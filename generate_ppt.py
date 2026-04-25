import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # ── Slide 1: Title Slide ───────────────────────────────────────────────
    slide_layout = prs.slide_layouts[0] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "ColonoSense QA Technical Overview"
    subtitle.text = "LLM Integration, Resource Requirements, and Evaluation Methodology\nPatient 4 Evaluation"
    
    # ── Slide 2: Technology Stack & LLM Usage ──────────────────────────────
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Technology Stack & LLM Utilization"
    
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "ColonoSense LLM Infrastructure"
    
    p = tf.add_paragraph()
    p.text = "Primary Model: OpenAI GPT-4o-mini (or Local Llama 3.1 70B via DGX)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Reasoning: Chosen for the optimal balance of fast token generation, deep clinical reasoning, and high template adherence."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Resource Requirements (Per Clinical Question):"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Context Input (Core RAG + Guard RAG): ~2,500 tokens"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Generation Output: ~400 tokens"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Total Footprint: ~2,900 tokens per query."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "For 18 evaluation categories, it takes ~52,200 tokens total per patient report."
    p.level = 1

    # ── Slide 3: RAG Template Matching Strategy ────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Response Template Matching Strategy"
    
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "How we enforce strict output templates without hallucination:"
    
    p = tf.add_paragraph()
    p.text = "1. Structured Patient Anchor:"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "We pre-calculate numeric variables (MES, Nancy, CRP) via Python and inject them as an 'Anchor Block'."
    p2.level = 2

    p = tf.add_paragraph()
    p.text = "2. Prompt Directives:"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "The LLM is explicitly commanded to copy values directly from the Anchor Block rather than inferring from narrative text."
    p2.level = 2

    p = tf.add_paragraph()
    p.text = "3. Fill-in-the-blank Validation:"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "By matching the exact formatting expected by the clinical grading rubric, we ensure 100% adherence to trial protocols."
    p2.level = 2

    # ── Slide 4: Evaluation Mathematical Equations ─────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "QA Evaluation Metrics & Mathematical Equations"
    
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Calculations used in the Evaluation Dashboard:"
    
    p = tf.add_paragraph()
    p.text = "Data Retrieval Accuracy: (Correctly Extracted Variables) / (Total Expected Variables)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Correctness: 1 - [(Hallucinations + Contradictions) / (Total Claims)]"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Concordance (Guideline Adherence): (Rules Adhered To) / (Total STRIDE-II Rules Applicable)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Completeness: (Output Fields Filled) / (Total Required Template Fields)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Helpfulness: Scored 1.0 if reasoning is logical and provides clear next clinical steps, else 0.0"
    p.level = 1

    # ── Slide 5: Report Snapshot Placeholder ───────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only layout
    title = slide.shapes.title
    title.text = "ColonoSense Patient Evaluation Report Snapshot"
    
    # Add a big placeholder shape for the screenshot
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4.5)
    
    shape = slide.shapes.add_shape(
        1, # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    
    tf = shape.text_frame
    tf.text = "[ PLEASE INSERT SCREENSHOT OF HTML REPORT HERE ]"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(220, 220, 220)
    
    # Save presentation
    prs.save("ColonoSense_QA_Technical_Overview.pptx")
    print("Presentation saved as ColonoSense_QA_Technical_Overview.pptx")

if __name__ == "__main__":
    create_presentation()
