import os
from html2image import Html2Image
from pptx import Presentation
from pptx.util import Inches, Pt

def add_screenshot_to_ppt():
    # 1. Take Screenshot
    hti = Html2Image()
    html_file = 'colonosense_report_patient4_20260423_2100.html'
    screenshot_file = 'report_screenshot.png'
    
    if os.path.exists(html_file):
        print(f"Taking screenshot of {html_file}...")
        hti.screenshot(html_file=html_file, save_as=screenshot_file, size=(1000, 800))
    else:
        print(f"Error: {html_file} not found.")
        return

    # 2. Add to PPT
    ppt_file = "ColonoSense_QA_Technical_Overview.pptx"
    prs = Presentation(ppt_file)
    
    # We want to add it to the last slide (Slide 5)
    slide = prs.slides[-1]
    
    # Remove the placeholder shape we added previously (it should be the last shape)
    # The title is shapes[0], the rectangle placeholder is shapes[1]
    if len(slide.shapes) > 1:
        sp = slide.shapes[1]._element
        sp.getparent().remove(sp)

    # Insert the screenshot image
    left = Inches(0.5)
    top = Inches(1.5)
    height = Inches(4.5)
    # let width auto-adjust to maintain aspect ratio
    pic = slide.shapes.add_picture(screenshot_file, left, top, height=height)
    
    # 3. Add descriptions pointing to the metrics
    desc_left = Inches(6.5)
    desc_top = Inches(1.5)
    desc_width = Inches(3.0)
    desc_height = Inches(4.5)
    
    textbox = slide.shapes.add_textbox(desc_left, desc_top, desc_width, desc_height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "Dashboard Metrics Explained:"
    p.font.bold = True
    p.font.size = Pt(16)
    
    p2 = tf.add_paragraph()
    p2.text = "• Data Retrieval (80%): Measures if the LLM correctly fetched numeric anchors without omitting data."
    p2.font.size = Pt(12)
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = "• Correctness (87.2%): Ensures zero hallucinations in the final diagnosis."
    p3.font.size = Pt(12)
    p3.space_before = Pt(6)
    
    p4 = tf.add_paragraph()
    p4.text = "• Concordance (77.2%): Assesses how strictly the AI adhered to the STRIDE-II clinical trial protocols."
    p4.font.size = Pt(12)
    p4.space_before = Pt(6)
    
    p5 = tf.add_paragraph()
    p5.text = "• Completeness (93.6%): Confirms all required fields in the forced template were populated."
    p5.font.size = Pt(12)
    p5.space_before = Pt(6)
    
    prs.save(ppt_file)
    print("Screenshot added to PPT successfully.")

if __name__ == "__main__":
    add_screenshot_to_ppt()
