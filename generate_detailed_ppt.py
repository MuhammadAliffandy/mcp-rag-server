import asyncio
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

async def generate_ppt():
    html_file = 'file://' + os.path.abspath('colonosense_report_patient4_20260423_2100.html')
    ppt_file = 'ColonoSense_QA_Technical_Overview_Detailed.pptx'
    
    # 1. Create PPT
    prs = Presentation()
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "ColonoSense Detailed QA Report"
    slide.placeholders[1].text = "Per-Question Breakdown & Evaluation Mathematics"
    
    # Metrics slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Evaluation Mathematics & Definitions"
    tf = slide.shapes.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "1. Data Retrieval: (Correct Anchors) / (Expected Anchors)"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "2. Correctness: 1 - [(Hallucinations + Contradictions) / (Total Facts)]"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "3. Concordance: (Matched STRIDE-II Rules) / (Applicable Rules)"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "4. Completeness: (Filled Template Fields) / (Total Required Fields)"
    p.font.size = Pt(16)

    # 2. Capture Screenshots using Playwright
    async with async_playwright() as p_w:
        browser = await p_w.chromium.launch(headless=True)
        page = await browser.new_page(device_scale_factor=2)
        
        try:
            await page.goto(html_file, wait_until='load')
        except Exception as e:
            print(f"Error opening HTML: {e}")
            await browser.close()
            return
            
        # Get all question elements (.qa)
        qa_elements = await page.locator('.qa').all()
        
        if not qa_elements:
            print("No .qa elements found in the HTML.")
            await browser.close()
            return
            
        for i, qa in enumerate(qa_elements):
            ss_path = f"ss_q_{i}.png"
            await qa.screenshot(path=ss_path)
            
            # Create a new slide for this question
            slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only
            slide.shapes.title.text = f"Question Evaluation {i+1}"
            
            # Insert screenshot
            left = Inches(0.5)
            top = Inches(1.5)
            # Add picture
            slide.shapes.add_picture(ss_path, left, top, width=Inches(9))
            
            # Add specific metric description
            desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
            dtf = desc_box.text_frame
            p = dtf.add_paragraph()
            p.text = "Calculation applied to this question:"
            p.font.bold = True
            
            p2 = dtf.add_paragraph()
            p2.text = "Retrieval = Exact match of Anchor Block. Completeness = No missing placeholders. Correctness = No hallucinations found. Concordance = Matches STRIDE-II targets."
            p2.font.size = Pt(12)
            
        await browser.close()

    # Save presentation
    prs.save(ppt_file)
    print(f"Presentation saved successfully as {ppt_file}")

if __name__ == "__main__":
    asyncio.run(generate_ppt())
